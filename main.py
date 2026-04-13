"""
NoteLM — NotebookLM-grade research assistant
FastAPI + Gemini backend  
Multi-notebook version with sharing
"""

from __future__ import annotations

import asyncio
import json
import os
import re
import secrets
import textwrap
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional

import httpx
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import networkx as nx
import numpy as np
from docx import Document
from gtts import gTTS
from fastapi import BackgroundTasks, FastAPI, File, Form, HTTPException, Query, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from fastapi.staticfiles import StaticFiles
import pdfplumber
from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from google import genai
import requests

MISTRAL_API_KEY = os.environ.get("MISTRAL_API_KEY", "")
MISTRAL_URL = "https://api.mistral.ai/v1/chat/completions"

# ──────────────────────────────────────────────────────────────────────────────
# Config
# ──────────────────────────────────────────────────────────────────────────────

BASE_DIR  = Path(__file__).parent
UPLOADS   = BASE_DIR / "uploads"
OUTPUTS   = BASE_DIR / "outputs"
STATIC    = BASE_DIR / "static"
DATA_DIR  = BASE_DIR / "data"

for d in (UPLOADS, OUTPUTS, STATIC, DATA_DIR):
    d.mkdir(exist_ok=True)

NOTEBOOKS_FILE   = DATA_DIR / "notebooks.json"
SOURCES_FILE     = DATA_DIR / "sources.json"
OLD_SOURCES_FILE = BASE_DIR / "sources.json"  # legacy migration

GEMINI_API_KEY = os.environ.get(
    "GEMINI_API_KEY", "AIzaSyBvnlnYaex2EecyzuJ9N3t1yInFII4U_zw"
)
_client = genai.Client(api_key=GEMINI_API_KEY)

# ──────────────────────────────────────────────────────────────────────────────
# Persistence
# ──────────────────────────────────────────────────────────────────────────────

def _load_json(path: Path, default=None):
    if path.exists():
        try:
            return json.loads(path.read_text())
        except Exception:
            pass
    return default if default is not None else {}

def _save_json(path: Path, data) -> None:
    path.write_text(json.dumps(data, indent=2, default=str))

notebooks_db: dict = {}
sources_db: dict   = {}

def _load_dbs():
    global notebooks_db, sources_db
    notebooks_db = _load_json(NOTEBOOKS_FILE, {})
    sources_db   = _load_json(SOURCES_FILE, {})

def _save_notebooks():
    _save_json(NOTEBOOKS_FILE, notebooks_db)

def _save_sources():
    _save_json(SOURCES_FILE, sources_db)

def _save_all():
    _save_notebooks()
    _save_sources()

# ── Migration from old flat sources.json ──────────────────────────────────────
def _migrate():
    if notebooks_db:
        return
    old = _load_json(OLD_SOURCES_FILE, {})
    if not old:
        return
    nid = "nb_" + uuid.uuid4().hex[:10]
    notebooks_db[nid] = {
        "id": nid,
        "name": "My Research",
        "description": "Migrated from previous version",
        "created_at": datetime.now().isoformat(),
        "sources": list(old.keys()),
        "share_token": None,
        "shared_at": None,
    }
    sources_db.update(old)
    _save_all()

# ──────────────────────────────────────────────────────────────────────────────
# Text extraction
# ──────────────────────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                pages.append(t)
    return "\n\n".join(pages)

def extract_docx(path: Path) -> str:
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def extract_txt(path: Path) -> str:
    return path.read_text(errors="ignore")

async def fetch_url(url: str) -> str:
    async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
        r = await client.get(url, headers={"User-Agent": "Mozilla/5.0"})
        r.raise_for_status()
        text = re.sub(r"<[^>]+>", " ", r.text)
        text = re.sub(r"\s+", " ", text)
        return text[:60000]

# ──────────────────────────────────────────────────────────────────────────────
# Gemini helpers
# ──────────────────────────────────────────────────────────────────────────────

def _gemini(prompt: str, quality: str = 'high') -> str:
    """quality='high' → gemini-3.1-pro-preview, 'fast' → gemini-2.5-flash"""
    model = 'gemini-3.1-pro-preview' if quality == 'high' else 'gemini-2.5-flash'
    try:
        resp = _client.models.generate_content(model=model, contents=prompt)
        return resp.text.strip()
    except Exception as e1:
        # Fallback to flash if pro fails
        if model != 'gemini-2.5-flash':
            try:
                resp = _client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
                return resp.text.strip()
            except Exception as gemini_err:
                pass
        else:
            gemini_err = e1
        if not MISTRAL_API_KEY:
            raise gemini_err
        headers = {"Authorization": f"Bearer {MISTRAL_API_KEY}", "Content-Type": "application/json"}
        payload = {"model": "mistral-medium-latest", "messages": [{"role": "user", "content": prompt}], "max_tokens": 8192}
        resp = requests.post(MISTRAL_URL, headers=headers, json=payload, timeout=90)
        resp.raise_for_status()
        return resp.json()["choices"][0]["message"]["content"].strip()

def _gemini_json(prompt: str, quality: str = 'high') -> dict | list:
    raw = _gemini(prompt, quality=quality)
    raw = re.sub(r"^```(?:json)?\s*", "", raw, flags=re.MULTILINE)
    raw = re.sub(r"```\s*$", "", raw, flags=re.MULTILINE)
    # Find JSON boundaries
    for start_char, end_char in [('{', '}'), ('[', ']')]:
        s = raw.find(start_char)
        e = raw.rfind(end_char)
        if s != -1 and e != -1:
            try:
                return json.loads(raw[s:e+1])
            except Exception:
                pass
    return json.loads(raw.strip())

# ──────────────────────────────────────────────────────────────────────────────
# Analysis pipeline
# ──────────────────────────────────────────────────────────────────────────────

def analyze(content: str, name: str, lang: str = 'ru') -> dict:
    lang_instruction = (
        "ПОЛНОСТЬЮ НА РУССКОМ ЯЗЫКЕ. Все поля JSON на русском. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "ENTIRELY IN ENGLISH. All JSON fields in English."
    )
    lang_hint = "Russian" if lang == "ru" else "English"

    prompt = f"""{lang_instruction}

You are a world-class research analyst. Perform DEEP, SUBSTANTIVE analysis of the document titled "{name}".

Be specific, not generic. Extract non-obvious insights. Avoid restating headings — find the real substance.

Requirements:
- summary: 5-6 detailed paragraphs covering ALL major points, arguments, and conclusions. Be thorough.
- key_points: 15+ specific, non-obvious insights. Each should be a complete thought with context, not a topic label.
- topics: 5-8 topics, each with 4-5 detailed sub-points that explain mechanisms, not just name them
- notable_quotes: 8-10 most impactful, interesting, or surprising quotes verbatim from the text
- entities: comprehensive extraction of all people, organizations, places, dates mentioned

Return ONLY a valid JSON object (no markdown fences) with exactly these fields:

{{
  "summary": "Comprehensive 5-6 paragraph executive summary covering all major themes",
  "key_points": ["specific insight with context 1", "specific insight 2", "...at least 15"],
  "topics": [
    {{
      "title": "Topic Name",
      "description": "2-3 sentence substantive description explaining what and why",
      "points": ["detailed sub-point 1", "detailed sub-point 2", "detailed sub-point 3", "detailed sub-point 4"]
    }}
  ],
  "notable_quotes": ["exact verbatim quote 1", "exact quote 2", "...8-10 quotes"],
  "entities": {{
    "people": ["name1", "..."],
    "organizations": ["org1", "..."],
    "places": ["place1", "..."],
    "dates": ["date1", "..."]
  }},
  "sentiment": "positive|negative|neutral|mixed",
  "complexity": "beginner|intermediate|advanced",
  "word_count": 1234,
  "language": "{lang_hint}"
}}

Document content (first 30000 chars):
{content[:30000]}"""

    try:
        return _gemini_json(prompt)
    except Exception as e:
        return {
            "summary": "Analysis failed. Content preview: " + content[:300],
            "key_points": [],
            "topics": [],
            "notable_quotes": [],
            "entities": {"people": [], "organizations": [], "places": [], "dates": []},
            "sentiment": "neutral",
            "complexity": "intermediate",
            "word_count": len(content.split()),
            "language": lang_hint,
            "_error": str(e),
        }

# ──────────────────────────────────────────────────────────────────────────────
# Output generators
# ──────────────────────────────────────────────────────────────────────────────

_DARK   = "#0F172A"
_MID    = "#334155"
_ACCENT = "#6366F1"
_LIGHT  = "#F8FAFC"
_MUTED  = "#94A3B8"


def _gen_rich_slide_content(source: dict, lang: str = 'ru') -> dict:
    """Use Gemini to generate rich, structured slide content for a presentation."""
    a = source.get("analysis") or {}
    content = source.get("content", "")[:20000]
    name = source.get("name", "Document")
    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. All JSON values in Russian."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE. All JSON values in English."
    )
    prompt = f"""{lang_instruction}

You are a world-class presentation designer. Create rich, compelling slide content for a professional presentation about "{name}".

Source analysis:
- Summary: {a.get('summary', '')[:1000]}
- Key points: {json.dumps(a.get('key_points', [])[:10], ensure_ascii=False)}
- Topics: {json.dumps([{{'title': t.get('title',''), 'description': t.get('description',''), 'points': t.get('points',[])[:4]}} for t in a.get('topics', [])[:6]], ensure_ascii=False)}
- Notable quotes: {json.dumps(a.get('notable_quotes', [])[:5], ensure_ascii=False)}
- Entities: {json.dumps(a.get('entities', {{}}), ensure_ascii=False)}
- Sentiment: {a.get('sentiment', 'neutral')}

Content excerpt: {content[:5000]}

Generate presentation content as JSON. Make it RICH, SPECIFIC, and NON-GENERIC.
Each slide should have a compelling headline, substantive body content, and concrete details.

Return ONLY valid JSON (no markdown fences):
{{
  "title": "Compelling presentation title",
  "subtitle": "Engaging one-line subtitle",
  "executive_summary": "2-3 sentence executive summary that hooks the audience",
  "key_insight": "Single most important insight from the document",
  "slides": [
    {{
      "title": "Slide title (punchy, max 8 words)",
      "type": "insight",
      "headline": "Bold claim or key finding",
      "body": "2-3 sentences of substantive explanation with specific details",
      "bullets": ["Specific point with data/evidence", "Another concrete finding", "Third key element"],
      "stat": "Key statistic or number if available",
      "accent_color": "#6366F1"
    }}
  ],
  "conclusion": "Memorable closing statement",
  "call_to_action": "What the audience should do or think about"
}}

Generate 6-8 slides covering: overview, key findings (2-3 slides), deep insights, implications, conclusion.
Each slide must have unique, specific content from the document — NO generic placeholders."""

    try:
        return _gemini_json(prompt, quality='high')
    except Exception:
        return {}


def gen_presentation(source: dict, out: Path, lang: str = 'ru') -> None:
    a = source.get("analysis") or {}
    name = source.get("name", "Untitled")

    # Generate rich AI content
    rich = _gen_rich_slide_content(source, lang)

    prs = PptxPresentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def _blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _parse_color(color: str):
        named = {"white": "#FFFFFF", "black": "#000000"}
        c = named.get(str(color).lower(), color).lstrip('#')
        if len(c) == 3:
            c = c[0]*2 + c[1]*2 + c[2]*2
        try:
            return int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        except Exception:
            return 99, 102, 241  # fallback accent

    def _bg(slide, color: str):
        bg = slide.background; fill = bg.fill; fill.solid()
        r, g, b = _parse_color(color)
        fill.fore_color.rgb = RGBColor(r, g, b)

    def _txt(slide, text, left, top, width, height,
             bold=False, size=18, color="#0F172A", align="LEFT", italic=False):
        from pptx.enum.text import PP_ALIGN
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = str(text)
        p.font.bold = bold; p.font.size = Pt(size); p.font.italic = italic
        r, g, b = _parse_color(color)
        p.font.color.rgb = RGBColor(r, g, b)
        p.alignment = {"LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER, "RIGHT": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        return box

    def _txt_multiline(slide, lines, left, top, width, height,
                       bold=False, size=14, color="#0F172A"):
        from pptx.enum.text import PP_ALIGN
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame; tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = str(line); p.font.size = Pt(size); p.font.bold = bold
            r, g, b = _parse_color(color)
            p.font.color.rgb = RGBColor(r, g, b)
            p.space_after = Pt(4)
        return box

    def _rect(slide, left, top, width, height, color):
        r, g, b = _parse_color(color)
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()
        return shape

    def _rounded_rect(slide, left, top, width, height, color, corner_size=0.1):
        from pptx.util import Emu
        r, g, b = _parse_color(color)
        # Use rounded rectangle shape (MSO_SHAPE.ROUNDED_RECTANGLE = 5)
        shape = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()
        return shape

    # ── Slide 1: Hero Title ───────────────────────────────────────────────────
    slide = _blank(); _bg(slide, _DARK)
    # Gradient-like accent bar at top
    _rect(slide, 0, 0, 13.33, 0.08, _ACCENT)
    # Large decorative circle
    _rounded_rect(slide, 10.5, 0.8, 2.2, 2.2, "#1E293B")
    _txt(slide, "✦", 11.1, 1.05, 1.2, 1.0, bold=True, size=48, color=_ACCENT, align="CENTER")

    title_text = rich.get("title") or name
    subtitle_text = rich.get("subtitle") or "AI-Generated Research Overview"
    insight_text = rich.get("executive_summary") or rich.get("key_insight") or a.get("summary", "")[:200]

    _txt(slide, title_text[:80], 0.7, 1.2, 9.5, 2.0, bold=True, size=36, color=_LIGHT)
    _rect(slide, 0.7, 3.3, 4.0, 0.04, _ACCENT)
    _txt(slide, subtitle_text[:100], 0.7, 3.55, 9.5, 0.6, size=18, color="#94A3B8")
    _txt(slide, textwrap.fill(insight_text[:250], 90), 0.7, 4.35, 10.5, 1.5, size=14, color="#64748B", italic=True)
    _txt(slide, datetime.now().strftime("%B %Y"), 0.7, 6.5, 4, 0.5, size=13, color=_MUTED)

    # ── Slide 2: Executive Summary ────────────────────────────────────────────
    slide = _blank(); _bg(slide, "#0F172A")
    _rect(slide, 0, 0, 0.06, 7.5, _ACCENT)
    _rect(slide, 0, 0, 13.33, 1.1, "#1E293B")
    summary_label = "Резюме" if lang == "ru" else "Executive Summary"
    _txt(slide, summary_label, 0.4, 0.2, 12, 0.7, bold=True, size=28, color=_LIGHT)

    summary_text = a.get("summary", "")
    # Split into paragraphs for better layout
    para1 = textwrap.fill(summary_text[:400], 95)
    para2 = textwrap.fill(summary_text[400:750], 95) if len(summary_text) > 400 else ""

    _txt(slide, para1, 0.4, 1.3, 12.4, 2.2, size=16, color="#CBD5E1")
    if para2:
        _txt(slide, para2, 0.4, 3.6, 12.4, 1.8, size=16, color="#94A3B8")

    key_insight = rich.get("key_insight") or ""
    if key_insight:
        _rect(slide, 0.4, 5.6, 12.5, 1.4, "#1E3A5F")
        _txt(slide, "💡 " + key_insight[:200], 0.7, 5.75, 12.0, 1.1, size=15, color="#93C5FD", bold=True)

    # ── AI-Generated Rich Slides ──────────────────────────────────────────────
    ai_slides = rich.get("slides", [])

    if ai_slides:
        for slide_data in ai_slides[:8]:
            slide = _blank(); _bg(slide, _LIGHT)
            slide_accent = slide_data.get("accent_color", _ACCENT)
            try:
                _parse_color(slide_accent)
            except Exception:
                slide_accent = _ACCENT

            # Header bar
            _rect(slide, 0, 0, 13.33, 1.15, slide_accent)
            slide_title = slide_data.get("title", "")[:70]
            _txt(slide, slide_title, 0.5, 0.15, 12.3, 0.85, bold=True, size=26, color="white")

            # Headline - big bold claim
            headline = slide_data.get("headline", "")
            if headline:
                _txt(slide, headline[:120], 0.5, 1.3, 12.3, 0.75, bold=True, size=20, color=_DARK)
                _rect(slide, 0.5, 2.1, 12.3, 0.03, "#E2E8F0")

            # Body text
            body = slide_data.get("body", "")
            body_y = 2.25 if headline else 1.35
            if body:
                _txt(slide, textwrap.fill(body[:350], 95), 0.5, body_y, 12.3, 1.4, size=15, color=_MID)

            # Bullet points
            bullets = slide_data.get("bullets", [])
            bullet_y = body_y + (1.6 if body else 0.2)
            for bi, bullet in enumerate(bullets[:5]):
                y = bullet_y + bi * 0.7
                if y > 7.0:
                    break
                _rounded_rect(slide, 0.5, y + 0.05, 0.28, 0.28, slide_accent)
                _txt(slide, "→", 0.5, y + 0.02, 0.28, 0.3, bold=True, size=11, color="white", align="CENTER")
                _txt(slide, str(bullet)[:140], 0.9, y, 11.8, 0.55, size=14, color=_DARK)

            # Stat callout (if available and there's space)
            stat = slide_data.get("stat", "")
            if stat and not bullets:
                _rounded_rect(slide, 4.5, 5.2, 4.3, 1.5, slide_accent)
                _txt(slide, str(stat)[:60], 4.7, 5.3, 3.9, 1.1, bold=True, size=22, color="white", align="CENTER")

    else:
        # Fallback: use analysis data with improved layout
        kps = a.get("key_points", [])
        if kps:
            slide = _blank(); _bg(slide, _LIGHT)
            _rect(slide, 0, 0, 13.33, 1.15, "#4F46E5")
            kp_label = "Ключевые моменты" if lang == "ru" else "Key Points"
            _txt(slide, kp_label, 0.5, 0.2, 12, 0.75, bold=True, size=28, color="white")
            for i, pt in enumerate(kps[:7]):
                y = 1.35 + i * 0.78
                _rounded_rect(slide, 0.5, y, 0.35, 0.35, _ACCENT)
                _txt(slide, str(i+1), 0.52, y+0.01, 0.3, 0.33, bold=True, size=13, color="white", align="CENTER")
                _txt(slide, str(pt)[:160], 1.0, y+0.01, 11.8, 0.55, size=15, color=_DARK)

        for topic in a.get("topics", [])[:6]:
            slide = _blank(); _bg(slide, _LIGHT)
            _rect(slide, 0, 0, 0.18, 7.5, _ACCENT)
            _txt(slide, topic.get("title", "")[:70], 0.45, 0.3, 12.4, 0.9, bold=True, size=26, color=_DARK)
            _rect(slide, 0.45, 1.25, 12.4, 0.04, "#E2E8F0")
            _txt(slide, topic.get("description", "")[:250], 0.45, 1.45, 12.2, 1.2, size=16, color=_MID)
            for j, pt in enumerate(topic.get("points", [])[:5]):
                _rect(slide, 0.45, 2.85+j*0.72, 0.06, 0.32, _ACCENT)
                _txt(slide, str(pt)[:160], 0.7, 2.82+j*0.72, 12.1, 0.55, size=15, color=_DARK)

    # ── Entities slide ────────────────────────────────────────────────────────
    ents = a.get("entities", {})
    if any(ents.values()):
        slide = _blank(); _bg(slide, _DARK)
        _rect(slide, 0, 0, 13.33, 1.15, "#1E293B")
        ent_label = "Ключевые сущности" if lang == "ru" else "Key Entities"
        _txt(slide, ent_label, 0.5, 0.2, 12, 0.75, bold=True, size=28, color=_LIGHT)
        for ci, (label_text, items, col) in enumerate([
            ("People" if lang == "en" else "Люди", ents.get("people", []), "#818CF8"),
            ("Organizations" if lang == "en" else "Организации", ents.get("organizations", []), "#34D399"),
            ("Places" if lang == "en" else "Места", ents.get("places", []), "#FB923C"),
            ("Dates" if lang == "en" else "Даты", ents.get("dates", []), "#F472B6"),
        ]):
            x = 0.4 + ci * 3.2
            _rounded_rect(slide, x, 1.25, 3.0, 0.52, col)
            _txt(slide, label_text, x+0.12, 1.32, 2.75, 0.38, bold=True, size=14, color="white")
            items_text = "\n".join(f"• {it}" for it in items[:7]) or "—"
            _txt(slide, items_text, x+0.12, 1.9, 2.85, 4.9, size=13, color=_MUTED)

    # ── Quotes slide ──────────────────────────────────────────────────────────
    quotes = a.get("notable_quotes", [])
    if quotes:
        slide = _blank(); _bg(slide, "#0F172A")
        _rect(slide, 0, 0, 13.33, 0.06, _ACCENT)
        quotes_label = "Примечательные цитаты" if lang == "ru" else "Notable Quotes"
        _txt(slide, quotes_label, 0.6, 0.25, 12, 0.8, bold=True, size=28, color=_LIGHT)
        for qi, q in enumerate(quotes[:3]):
            y = 1.25 + qi * 1.9
            _rect(slide, 0.4, y, 0.06, 1.55, _ACCENT)
            _rounded_rect(slide, 0.6, y, 12.3, 1.55, "#1E293B")
            _txt(slide, f'"{str(q)[:300]}"', 0.9, y + 0.15, 11.8, 1.2, size=15, color="#CBD5E1", italic=True)

    # ── Conclusion slide ──────────────────────────────────────────────────────
    slide = _blank(); _bg(slide, _DARK)
    _rect(slide, 0, 7.3, 13.33, 0.2, _ACCENT)
    conclusion_label = "Заключение" if lang == "ru" else "Conclusion"
    _txt(slide, conclusion_label, 0.7, 0.9, 11.9, 0.8, bold=True, size=32, color=_LIGHT, align="CENTER")
    _rect(slide, 4.5, 1.8, 4.33, 0.04, _ACCENT)

    conclusion = rich.get("conclusion") or a.get("summary", "")[:300]
    cta = rich.get("call_to_action", "")

    _txt(slide, textwrap.fill(str(conclusion)[:400], 80), 1.2, 2.1, 10.9, 2.5, size=18, color="#CBD5E1", align="CENTER")
    if cta:
        _rounded_rect(slide, 2.5, 5.0, 8.33, 1.2, "#1E293B")
        _txt(slide, str(cta)[:180], 2.8, 5.15, 7.7, 0.9, size=15, color=_ACCENT, align="CENTER", bold=True)

    prs.save(str(out))


def gen_audio(source: dict, out: Path, lang: str = 'ru') -> None:
    """Generate a natural, conversational audio overview using Gemini for the script."""
    a = source.get("analysis") or {}
    content = source.get("content", "")[:15000]
    name = source.get("name", "this document")

    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE."
    )
    lang_label = "Russian" if lang == "ru" else "English"

    prompt = f"""{lang_instruction}

You are a professional podcast host. Write a natural, conversational audio script
reviewing the document "{name}". The script should be 3-5 minutes long when read aloud
(approximately 500-750 words). Do NOT include any stage directions, sound effects,
or formatting marks. Write only the spoken words.

Start with a friendly introduction, cover the main topics and key insights,
include interesting details and context, and end with a thoughtful conclusion.
Use natural speech patterns, transitions like "Now let's talk about...", "What's interesting here is...",
"One key point to understand is...". Make it engaging and informative.

Generate ONLY the script text in {lang_label}, nothing else.

Document summary: {a.get('summary', '')}
Key points: {', '.join(a.get('key_points', [])[:8])}
Main topics: {', '.join(t.get('title','') for t in a.get('topics',[])[:5])}
Content excerpt: {content[:5000]}"""

    try:
        script = _gemini(prompt)
    except Exception:
        # Fallback script
        kps = a.get("key_points", [])
        topics = a.get("topics", [])
        if lang == "ru":
            script = (
                f"Добро пожаловать в аудиообзор документа «{name}». "
                f"{a.get('summary', 'Этот документ содержит важную информацию.')} "
            )
            if kps:
                script += "Среди ключевых выводов стоит отметить: " + ". ".join(kps[:6]) + ". "
            if topics:
                script += "Документ охватывает следующие темы: " + ", ".join(
                    t["title"] for t in topics[:5]) + ". "
            script += "Надеемся, этот обзор был полезен. Спасибо за внимание."
        else:
            script = (
                f"Welcome to the audio overview of '{name}'. "
                f"{a.get('summary', 'This document contains important information.')} "
            )
            if kps:
                script += "Key takeaways include: " + ". ".join(kps[:6]) + ". "
            if topics:
                script += "The main topics covered are: " + ", ".join(
                    t["title"] for t in topics[:5]) + ". "
            script += "We hope this overview was helpful. Thank you for listening."

    tts_lang = "ru" if lang == "ru" else "en"
    tts = gTTS(text=script, lang=tts_lang, slow=False)
    tts.save(str(out))


def build_infographic_prompt(source: dict, lang: str = 'ru') -> str:
    a = source.get('analysis') or {}
    name = source.get('name', 'Document')
    summary = a.get('summary', '')[:500]
    key_points = a.get('key_points', [])[:5]
    topics = [t['title'] for t in a.get('topics', [])[:4]]
    sentiment = a.get('sentiment', 'neutral')

    lang_instruction = "All text in the infographic must be in Russian language." if lang == 'ru' else "All text in the infographic must be in English language."

    points_text = '; '.join(key_points) if key_points else 'key insights from the document'
    topics_text = ', '.join(topics) if topics else 'main topics'

    prompt = f"""Create a modern, beautiful, professional infographic about: "{name}".
{lang_instruction}

Design requirements:
- Clean white or light background
- Modern flat design with vibrant accent colors
- Clear typography, large readable text
- Visual hierarchy: title at top, key sections clearly separated
- Include: document title, 3-5 key points as bullet items, topic bubbles or icons, sentiment indicator
- Style: similar to Canva or Google Slides infographic templates
- Dimensions: landscape orientation, 16:9 ratio

Content to visualize:
- Title: {name}
- Summary theme: {summary[:200]}
- Key points: {points_text}
- Main topics: {topics_text}
- Sentiment: {sentiment}

Make it visually stunning with icons, charts, or visual elements. Professional quality."""
    return prompt


def _gen_infographic_matplotlib(source: dict, out: Path) -> None:
    """Fallback matplotlib infographic generation."""
    a = source.get("analysis") or {}
    topics = a.get("topics", [])
    kps    = a.get("key_points", [])
    ents   = a.get("entities", {})

    fig = plt.figure(figsize=(16, 10), facecolor=_DARK)
    fig.patch.set_facecolor(_DARK)
    gs = fig.add_gridspec(2, 3, hspace=0.45, wspace=0.4,
                          left=0.06, right=0.97, top=0.88, bottom=0.08)
    fig.text(0.5, 0.95, source.get("name", "Document Overview"),
             ha="center", va="top", fontsize=20, fontweight="bold", color="white")
    fig.text(0.5, 0.91, "AI-Generated Infographic · NoteLM",
             ha="center", va="top", fontsize=11, color=_MUTED)

    accent_colors = ["#6366F1","#8B5CF6","#EC4899","#F59E0B","#10B981","#3B82F6","#EF4444"]

    ax1 = fig.add_subplot(gs[0, :2]); ax1.set_facecolor("#1E293B")
    if topics:
        labels = [t["title"][:22] for t in topics[:7]]
        vals   = [max(len(t.get("points",[])), 1) for t in topics[:7]]
        bars   = ax1.barh(labels[::-1], vals[::-1], color=accent_colors[:len(labels)], height=0.55)
        for bar, val in zip(bars, vals[::-1]):
            ax1.text(bar.get_width()+0.05, bar.get_y()+bar.get_height()/2,
                     str(val), va="center", color="white", fontsize=10)
        ax1.set_xlabel("Sub-points", color=_MUTED, fontsize=10)
        ax1.set_title("Topics by Depth", color="white", fontsize=13, pad=8)
        ax1.set_xlim(0, max(max(vals)+2, 5))
    else:
        ax1.text(0.5, 0.5, "No topics extracted", ha="center", color=_MUTED,
                 transform=ax1.transAxes, fontsize=12)
        ax1.set_title("Topics", color="white", fontsize=13)
    ax1.tick_params(colors="white", labelsize=9)
    ax1.spines[:].set_color("#334155")

    ax2 = fig.add_subplot(gs[0, 2]); ax2.set_facecolor("#1E293B")
    sentiment = a.get("sentiment", "neutral")
    sent_map = {
        "positive": ([0.8,0.1,0.1], ["#10B981","#334155","#EF4444"], "Mostly Positive"),
        "negative": ([0.1,0.1,0.8], ["#10B981","#334155","#EF4444"], "Mostly Negative"),
        "mixed":    ([0.4,0.2,0.4], ["#10B981","#334155","#EF4444"], "Mixed"),
        "neutral":  ([0.1,0.8,0.1], ["#10B981","#F59E0B","#EF4444"], "Neutral"),
    }
    sizes, colors, label = sent_map.get(sentiment, sent_map["neutral"])
    ax2.pie(sizes, colors=colors, startangle=90, wedgeprops=dict(width=0.5, edgecolor=_DARK))
    ax2.text(0, 0, label, ha="center", va="center", color="white", fontsize=9, fontweight="bold")
    ax2.set_title("Sentiment", color="white", fontsize=13, pad=8)

    ax3 = fig.add_subplot(gs[1, 0]); ax3.set_facecolor("#1E293B")
    ent_labels = ["People","Orgs","Places","Dates"]
    ent_vals   = [len(ents.get("people",[])), len(ents.get("organizations",[])),
                  len(ents.get("places",[])), len(ents.get("dates",[]))]
    ent_cols = ["#818CF8","#34D399","#FB923C","#F472B6"]
    ax3.bar(ent_labels, ent_vals, color=ent_cols, width=0.55)
    ax3.set_title("Entities Detected", color="white", fontsize=13, pad=8)
    ax3.tick_params(colors="white", labelsize=9); ax3.spines[:].set_color("#334155")
    for i, v in enumerate(ent_vals):
        ax3.text(i, v+0.1, str(v), ha="center", color="white", fontsize=10)
    ax3.set_ylim(0, max(max(ent_vals)+2, 5))

    ax4 = fig.add_subplot(gs[1, 1]); ax4.set_facecolor("#1E293B")
    if kps:
        lens = [len(k.split()) for k in kps[:10]]
        ax4.barh(range(len(lens)), lens,
                 color=["#6366F1" if l > 10 else "#8B5CF6" for l in lens], height=0.6)
        ax4.set_yticks(range(len(lens)))
        ax4.set_yticklabels([f"#{i+1}" for i in range(len(lens))], color="white", fontsize=9)
        ax4.set_xlabel("Words", color=_MUTED, fontsize=9)
    ax4.set_title("Key Point Lengths", color="white", fontsize=13, pad=8)
    ax4.tick_params(axis="x", colors="white", labelsize=9)
    ax4.spines[:].set_color("#334155")

    ax5 = fig.add_subplot(gs[1, 2]); ax5.set_facecolor("#1E293B"); ax5.axis("off")
    stats = [("Word Count", f"{a.get('word_count',0):,}"), ("Topics", str(len(topics))),
             ("Key Points", str(len(kps))), ("Complexity", a.get("complexity","—").title()),
             ("Language", a.get("language","—"))]
    ax5.set_title("Document Stats", color="white", fontsize=13, pad=8)
    for si, (lbl, val) in enumerate(stats):
        y = 0.85 - si * 0.18
        ax5.text(0.05, y, lbl, transform=ax5.transAxes, color=_MUTED, fontsize=11)
        ax5.text(0.95, y, val, transform=ax5.transAxes, color="white", fontsize=12,
                 fontweight="bold", ha="right")
        ax5.plot([0.05,0.95], [y-0.04,y-0.04], color="#334155", linewidth=0.5,
                 transform=ax5.transAxes, clip_on=False)

    plt.savefig(str(out), dpi=150, bbox_inches="tight", facecolor=_DARK, edgecolor="none")
    plt.close(fig)


import logging

def gen_infographic(source: dict, out: Path, lang: str = 'ru') -> None:
    """Clean professional PIL infographic — pure design, no AI background."""
    import logging
    from PIL import Image as PILImage, ImageDraw, ImageFont

    a    = source.get("analysis") or {}
    deep = source.get("deep_analysis") or {}
    name = source.get("name", "Document")
    pad  = 44
    W    = 1200

    ACCENTS = ['#6366F1','#10B981','#F59E0B','#0EA5E9','#EF4444','#8B5CF6','#EC4899']

    def h(hx):
        c = hx.lstrip('#')
        return tuple(int(c[i:i+2],16) for i in (0,2,4))

    def lf(bold=False, size=18):
        for p in [
            f'/usr/share/fonts/truetype/dejavu/DejaVuSans{"-Bold" if bold else ""}.ttf',
            f'/usr/share/fonts/truetype/noto/NotoSans{"-Bold" if bold else ""}.ttf',
        ]:
            if os.path.exists(p):
                try: return ImageFont.truetype(p, size)
                except: pass
        return ImageFont.load_default()

    F = {
        'h1':lf(True,40),'h2':lf(True,26),'h3':lf(True,20),
        'body':lf(False,17),'sm':lf(False,14),'num':lf(True,16),
    }

    def wrap(txt, font, maxw, draw):
        words = str(txt).split(); lines, cur = [], []
        for w in words:
            test = ' '.join(cur+[w])
            if draw.textbbox((0,0),test,font=font)[2] <= maxw: cur.append(w)
            else:
                if cur: lines.append(' '.join(cur))
                cur = [w]
        if cur: lines.append(' '.join(cur))
        return lines or ['']

    def rr(draw, x1,y1,x2,y2,r=10,fill=None):
        if not fill: return
        draw.rectangle([x1+r,y1,x2-r,y2],fill=fill)
        draw.rectangle([x1,y1+r,x2,y2-r],fill=fill)
        for cx,cy in [(x1,y1),(x2-2*r,y1),(x1,y2-2*r),(x2-2*r,y2-2*r)]:
            draw.ellipse([cx,cy,cx+2*r,cy+2*r],fill=fill)

    summary   = str(a.get('summary',''))[:500]
    kps       = a.get('key_points',[])[:7]
    topics    = a.get('topics',[])[:6]
    quotes    = a.get('notable_quotes',[])
    core      = str(deep.get('core_thesis',''))[:280]
    wcount    = a.get('word_count', len(summary.split()))
    sentiment = a.get('sentiment','neutral')
    complexity= a.get('complexity','intermediate')

    dummy = PILImage.new('RGB',(W,10)); dd = ImageDraw.Draw(dummy)

    def sec_h(k):
        if k=='header': return 185
        if k=='summary': return len(wrap(summary[:450],F['body'],W-2*pad-40,dd))*26+110
        if k=='kps':
            h_=70
            for kp in kps: h_+=max(len(wrap(kp[:200],F['body'],W-2*pad-88,dd)),1)*23+28
            return h_+20
        if k=='topics':
            h_=70; cw=(W-2*pad-16)//2
            for t in topics:
                h_+=max(len(wrap(t.get('title','')[:60],F['h3'],W-2*pad-60,dd)),1)*27
                h_+=max(len(wrap(t.get('description','')[:160],F['sm'],W-2*pad-60,dd)),1)*20+28
            return h_+20
        if k=='quote' and quotes:
            return len(wrap(f'"{quotes[0][:280]}"',F['body'],W-2*pad-60,dd))*28+80
        if k=='stats': return 110
        if k=='insight' and core:
            return len(wrap(core,F['body'],W-2*pad-40,dd))*26+110
        return 0

    secs = ['header','summary','kps']
    if topics: secs.append('topics')
    if quotes: secs.append('quote')
    secs.append('stats')
    if core: secs.append('insight')

    H = sum(sec_h(s) for s in secs)+60; H = max(H,1000)
    img = PILImage.new('RGB',(W,H),(255,255,255))
    draw = ImageDraw.Draw(img)
    y = 0

    for sec in secs:
        if sec=='header':
            draw.rectangle([0,0,W,180],fill=h('#0F172A'))
            draw.rectangle([0,176,W,180],fill=h('#6366F1'))
            clean_name = name.replace('_',' ').replace('.pdf','').replace('.docx','').replace('.txt','').replace('.md','')
            tl = wrap(clean_name[:80],F['h1'],W-2*pad,draw)
            draw.text((pad,32),tl[0],font=F['h1'],fill=(255,255,255))
            if len(tl)>1: draw.text((pad,82),tl[1],font=F['h1'],fill=(255,255,255))
            sub = ('Анализ документа' if lang=='ru' else 'Document Analysis')+' · Knowledge Studio'
            draw.text((pad,148),sub,font=F['sm'],fill=h('#94A3B8'))
            y = 200

        elif sec=='summary':
            sh = sec_h('summary')
            rr(draw,pad,y,W-pad,y+sh-10,14,fill=h('#F8FAFC'))
            draw.rectangle([pad,y,pad+6,y+sh-10],fill=h('#6366F1'))
            lbl = 'Краткое содержание' if lang=='ru' else 'Summary'
            draw.text((pad+20,y+14),lbl,font=F['h3'],fill=h('#6366F1'))
            ty=y+46
            for line in wrap(summary[:450],F['body'],W-2*pad-40,draw)[:13]:
                draw.text((pad+20,ty),line,font=F['body'],fill=h('#1E293B')); ty+=26
            y+=sh+10

        elif sec=='kps':
            lbl = 'Ключевые инсайты' if lang=='ru' else 'Key Insights'
            draw.text((pad,y),lbl,font=F['h2'],fill=h('#0F172A')); y+=36
            draw.rectangle([pad,y,pad+50,y+2],fill=h('#6366F1')); y+=14
            for i,kp in enumerate(kps):
                acc=h(ACCENTS[i%len(ACCENTS)])
                bg=h('#FAFAFA') if i%2==0 else h('#F1F5F9')
                lines=wrap(kp[:200],F['body'],W-2*pad-88,draw)
                ch=max(len(lines),1)*23+22
                rr(draw,pad,y,W-pad,y+ch,8,fill=bg)
                draw.rectangle([pad,y,pad+5,y+ch],fill=acc)
                cx2,cy2,r2=pad+26,y+ch//2,12
                draw.ellipse([cx2-r2,cy2-r2,cx2+r2,cy2+r2],fill=acc)
                nt=str(i+1); nb=draw.textbbox((0,0),nt,font=F['num'])
                draw.text((cx2-(nb[2]-nb[0])//2,cy2-8),nt,font=F['num'],fill=(255,255,255))
                for li,line in enumerate(lines[:3]):
                    draw.text((pad+48,y+7+li*23),line,font=F['body'],fill=h('#1E293B'))
                y+=ch+8
            y+=16

        elif sec=='topics':
            lbl = 'Основные темы' if lang=='ru' else 'Main Topics'
            draw.text((pad,y),lbl,font=F['h2'],fill=h('#0F172A')); y+=36
            draw.rectangle([pad,y,pad+50,y+2],fill=h('#10B981')); y+=14
            # Single-column for better readability — no text overflow
            for i,t in enumerate(topics):
                acc=h(ACCENTS[(i+2)%len(ACCENTS)])
                ttl = t.get('title','')[:70]
                desc = t.get('description','')[:160]
                tlines=wrap(ttl,F['h3'],W-2*pad-60,draw)
                dlines=wrap(desc,F['sm'],W-2*pad-60,draw)
                th=len(tlines)*27+len(dlines[:3])*20+36
                th=max(th,64)
                rr(draw,pad,y,W-pad,y+th,8,fill=h('#F8FAFC'))
                draw.rectangle([pad,y,pad+5,y+th],fill=acc)
                # Emoji number
                draw.ellipse([pad+10,y+th//2-12,pad+34,y+th//2+12],fill=acc)
                nt=str(i+1); nb=draw.textbbox((0,0),nt,font=F['num'])
                draw.text((pad+22-(nb[2]-nb[0])//2,y+th//2-8),nt,font=F['num'],fill=(255,255,255))
                ty2=y+8
                for tl2 in tlines[:2]:
                    draw.text((pad+46,ty2),tl2,font=F['h3'],fill=h('#0F172A')); ty2+=27
                for dl in dlines[:3]:
                    draw.text((pad+46,ty2),dl,font=F['sm'],fill=h('#6B7280')); ty2+=20
                y+=th+10
            y+=10

        elif sec=='quote':
            q=quotes[0][:280]
            qlines=wrap(f'"{q}"',F['body'],W-2*pad-60,draw)
            ch=len(qlines)*28+70
            rr(draw,pad,y,W-pad,y+ch,12,fill=h('#3730A3'))
            draw.text((pad+18,y+10),'❝',font=F['h2'],fill=h('#A5B4FC'))
            for li,line in enumerate(qlines):
                draw.text((pad+22,y+50+li*28),line,font=F['body'],fill=(220,225,255))
            y+=ch+16

        elif sec=='stats':
            sent_icon={'positive':'↑','negative':'↓','neutral':'→','mixed':'⟷'}.get(sentiment,'→')
            sent_lbl=('Позитивный' if sentiment=='positive' else 'Негативный' if sentiment=='negative' else 'Нейтральный') if lang=='ru' else sentiment.title()
            stats=[
                (str(wcount),'слов' if lang=='ru' else 'words','#6366F1'),
                (str(len(kps)),'инсайтов' if lang=='ru' else 'insights','#10B981'),
                (sent_icon+' '+sent_lbl,'тональность' if lang=='ru' else 'sentiment','#F59E0B'),
                (complexity.title(),'уровень' if lang=='ru' else 'level','#0EA5E9'),
            ]
            sw=(W-2*pad-16*3)//4
            for si,(val,lbl2,col2) in enumerate(stats):
                sx=pad+si*(sw+16)
                rr(draw,sx,y,sx+sw,y+90,10,fill=h(col2))
                vb=draw.textbbox((0,0),val[:12],font=F['h2']); vw=vb[2]-vb[0]
                draw.text((sx+(sw-vw)//2,y+12),val[:12],font=F['h2'],fill=(255,255,255))
                lb=draw.textbbox((0,0),lbl2[:16],font=F['sm']); lw=lb[2]-lb[0]
                draw.text((sx+(sw-lw)//2,y+56),lbl2[:16],font=F['sm'],fill=(220,230,255))
            y+=106

        elif sec=='insight':
            ilines=wrap(core,F['body'],W-2*pad-40,draw)
            ch=len(ilines)*26+90
            rr(draw,pad,y,W-pad,y+ch,12,fill=h('#0F172A'))
            draw.rectangle([pad,y,pad+6,y+ch],fill=h('#F59E0B'))
            lbl = '💡 Главный вывод' if lang=='ru' else '💡 Core Insight'
            draw.text((pad+20,y+14),lbl,font=F['h3'],fill=h('#F59E0B'))
            for li,line in enumerate(ilines):
                draw.text((pad+20,y+52+li*26),line,font=F['body'],fill=(200,210,230))
            y+=ch+16

    draw.rectangle([0,H-32,W,H],fill=h('#F1F5F9'))
    draw.text((pad,H-22),'Knowledge Studio',font=F['sm'],fill=h('#9CA3AF'))
    img.save(str(out),'PNG',dpi=(150,150))
    logging.info(f"Infographic saved: {out.stat().st_size//1024}KB")


def gen_mindmap(source: dict, out: Path) -> None:
    a = source.get("analysis") or {}
    topics = a.get("topics", [])
    name   = source.get("name", "Source")[:30]

    G = nx.Graph()
    center = name; G.add_node(center, level=0)
    topic_nodes = []; sub_nodes = []

    for topic in topics[:8]:
        tname = topic["title"][:28]; G.add_node(tname, level=1)
        G.add_edge(center, tname); topic_nodes.append(tname)
        for pt in topic.get("points", [])[:4]:
            pname = pt[:30]
            if pname in G.nodes: pname = pname + " "
            G.add_node(pname, level=2); G.add_edge(tname, pname); sub_nodes.append(pname)

    fig, ax = plt.subplots(figsize=(18, 12), facecolor=_DARK)
    ax.set_facecolor(_DARK); ax.axis("off")
    pos = nx.spring_layout(G, k=2.8, seed=42, iterations=80)
    nx.draw_networkx_edges(G, pos, ax=ax, edge_color="#334155", width=1.5, alpha=0.7)
    nx.draw_networkx_nodes(G, pos, ax=ax, nodelist=[center],
                           node_color=_ACCENT, node_size=3200, alpha=0.95)
    if topic_nodes:
        nx.draw_networkx_nodes(G, pos, ax=ax, nodelist=topic_nodes,
                               node_color="#4F46E5", node_size=1600, alpha=0.9)
    if sub_nodes:
        nx.draw_networkx_nodes(G, pos, ax=ax, nodelist=sub_nodes,
                               node_color="#1E293B", node_size=700, alpha=0.85)
    nx.draw_networkx_labels(G, pos, labels={center: textwrap.fill(center, 14)}, ax=ax,
                            font_size=11, font_color="white", font_weight="bold")
    nx.draw_networkx_labels(G, pos, labels={n: textwrap.fill(n, 14) for n in topic_nodes},
                            ax=ax, font_size=9, font_color="#C7D2FE")
    nx.draw_networkx_labels(G, pos, labels={n: textwrap.fill(n.strip(), 16) for n in sub_nodes},
                            ax=ax, font_size=7.5, font_color=_MUTED)
    fig.text(0.5, 0.97, f"Mind Map — {name}", ha="center",
             fontsize=16, fontweight="bold", color="white")
    fig.text(0.5, 0.935, "Generated by NoteLM", ha="center", fontsize=10, color=_MUTED)
    plt.tight_layout()
    plt.savefig(str(out), dpi=150, bbox_inches="tight", facecolor=_DARK, edgecolor="none")
    plt.close(fig)


def gen_study_guide(source: dict, lang: str = 'ru') -> str:
    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE."
    )
    a = source.get("analysis") or {}
    content = source.get("content", "")[:20000]
    lang_label = "Russian" if lang == "ru" else "English"
    prompt = f"""{lang_instruction}

Create a COMPREHENSIVE study guide for the document "{source.get('name','')}".
Generate ALL content in {lang_label} language.
Include 10+ review questions with detailed answers. Be thorough and educational.

Document summary: {a.get('summary', '')}

Full content excerpt:
{content}

The study guide MUST include all these sections:
1. **Overview** — What this document is about (2-3 paragraphs)
2. **Core Concepts** — All major concepts explained clearly with examples (at least 5)
3. **Key Terms & Definitions** — Important vocabulary (15+ terms, each with full definition)
4. **Concept Summary Table** — Two-column markdown table: Concept | Explanation (10+ rows)
5. **Critical Questions** — 10 thought-provoking questions with detailed answers (3-5 sentences each)
6. **Review Questions** — 10 factual Q&A pairs with thorough answers
7. **Practice Activities** — 3 hands-on exercises with step-by-step instructions

Format as clean markdown. Be thorough, detailed, and educational."""
    return _gemini(prompt, quality='high')


def gen_faq(source: dict, lang: str = 'ru') -> str:
    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE."
    )
    a = source.get("analysis") or {}
    content = source.get("content", "")[:20000]
    lang_label = "Russian" if lang == "ru" else "English"
    prompt = f"""{lang_instruction}

Generate 15-20 FAQ pairs for "{source.get('name','')}".
Generate ALL content in {lang_label} language.
Each answer MUST be 2-4 sentences minimum — substantive and informative.

Summary: {a.get('summary', '')}
Content: {content}

Create 15-20 natural FAQ pairs covering:
- Basic understanding questions (what, who, when, where)
- Clarification questions (how, why)
- Application questions (practical use cases)
- Edge case questions (limitations, exceptions)

Format as markdown:
## Q: [Question]
**A:** [Detailed answer, 2-4 sentences minimum]

Generate all 15-20 pairs. Do not skip any."""
    return _gemini(prompt, quality='high')


def gen_briefing(source: dict, lang: str = 'ru') -> str:
    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE."
    )
    a = source.get("analysis") or {}
    content = source.get("content", "")[:20000]
    lang_label = "Russian" if lang == "ru" else "English"
    prompt = f"""{lang_instruction}

Write a thorough executive briefing for "{source.get('name','')}".
Generate ALL content in {lang_label} language.
The briefing MUST be at least 500 words total. Be substantive and detailed.

# Executive Briefing: {source.get('name','')}

## Bottom Line Up Front (BLUF)
2-3 sentence critical takeaway.

## Situation
Current context and background (3-4 paragraphs, thorough).

## Key Findings
Numbered list of 7-10 most important findings, each with explanation (2-3 sentences per finding).

## Implications
What this means and why it matters (2-3 paragraphs with specific details).

## Recommendations
5-7 actionable recommendations, each with reasoning and expected outcome.

## Supporting Data
Key statistics, quotes, or evidence with context and interpretation.

## Risk Factors
4-6 potential concerns or caveats, each explained clearly.

Content: {content}
Summary: {a.get('summary', '')}"""
    return _gemini(prompt, quality='high')


def gen_timeline(source: dict, lang: str = 'ru') -> str:
    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE."
    )
    a = source.get("analysis") or {}
    content = source.get("content", "")[:20000]
    dates = a.get("entities", {}).get("dates", [])
    lang_label = "Russian" if lang == "ru" else "English"
    prompt = f"""{lang_instruction}

Extract and reconstruct a chronological timeline from "{source.get('name','')}".
Generate ALL content in {lang_label} language.

Known dates found: {', '.join(dates) if dates else 'scan content for any dates/periods'}.

Content: {content}

Format as a markdown timeline:
## Timeline

| Date/Period | Event/Development |
|------------|------------------|
| ...        | ...              |

Include ALL dates, periods, and sequential events you can identify.
After the table, add a **Narrative Summary** paragraph explaining the chronology in detail."""
    return _gemini(prompt, quality='high')


def gen_glossary(source: dict, lang: str = 'ru') -> str:
    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE."
    )
    content = source.get("content", "")[:20000]
    lang_label = "Russian" if lang == "ru" else "English"
    prompt = f"""{lang_instruction}

Extract all technical terms, jargon, and important concepts from "{source.get('name','')}".
Generate ALL content in {lang_label} language. Include 15-30 terms minimum.

Content: {content}

Create a comprehensive glossary:
# Glossary — {source.get('name','')}

For each term provide:
**Term** — Clear, concise definition in plain language. Note context and importance if relevant.

List terms alphabetically. Include at minimum 15-30 terms."""
    return _gemini(prompt, quality='high')


# ──────────────────────────────────────────────────────────────────────────────
# Deep analysis
# ──────────────────────────────────────────────────────────────────────────────

def deep_analyze(source: dict, lang: str = 'ru') -> dict:
    """Deep analysis: connections, semantics, logic, key conclusions."""
    content = source.get('content', '')
    name = source.get('name', '')
    lang_word = "RUSSIAN" if lang == "ru" else "ENGLISH"

    prompt = f"""RESPOND ENTIRELY IN {lang_word} LANGUAGE.

Perform DEEP ANALYSIS of this document: "{name}"

Analyze:
1. SEMANTIC STRUCTURE - how ideas build on each other
2. LOGICAL CHAIN - cause-effect relationships, argumentation
3. KEY INSIGHTS - non-obvious conclusions the author implies
4. CONTRADICTIONS - internal contradictions or tensions in the text
5. CORE THESIS - the single most important idea
6. EVIDENCE MAP - what arguments support what conclusions
7. IMPLICIT KNOWLEDGE - what the author assumes the reader knows
8. PRACTICAL APPLICATIONS - how to apply this knowledge

Content: {content[:25000]}

Return ONLY valid JSON (no markdown fences):
{{
  "core_thesis": "string",
  "semantic_structure": [{{"level": 1, "idea": "string", "builds_on": null}}],
  "logical_chains": [{{"premise": "string", "conclusion": "string", "strength": "strong|medium|weak"}}],
  "key_insights": [{{"insight": "string", "evidence": "string", "importance": "high|medium|low"}}],
  "contradictions": [{{"statement_a": "string", "statement_b": "string", "type": "internal|external"}}],
  "evidence_map": [{{"claim": "string", "supporting_evidence": ["string"]}}],
  "practical_applications": ["string"],
  "podcast_talking_points": ["key point for engaging podcast discussion"]
}}"""

    try:
        return _gemini_json(prompt)
    except Exception as e:
        import logging
        logging.warning(f"Deep analysis failed for {source.get('id','?')}: {e}")
        return {
            "core_thesis": "",
            "semantic_structure": [],
            "logical_chains": [],
            "key_insights": [],
            "contradictions": [],
            "evidence_map": [],
            "practical_applications": [],
            "podcast_talking_points": [],
            "_error": str(e),
        }


# ──────────────────────────────────────────────────────────────────────────────
# Book processing pipeline
# ──────────────────────────────────────────────────────────────────────────────

def process_book(source: dict, content: str, lang: str = 'ru') -> dict:
    """Extract TOC, split by chapters, generate summaries."""
    lang_instruction = (
        "RESPOND ENTIRELY IN RUSSIAN LANGUAGE. НЕ ИСПОЛЬЗУЙ АНГЛИЙСКИЙ ЯЗЫК."
        if lang == "ru" else "RESPOND ENTIRELY IN ENGLISH LANGUAGE."
    )

    # 1. Detect TOC with regex/heuristic
    chapters = []
    lines = content.split('\n')
    toc_patterns = [
        r'^(Chapter|Глава|Часть|Part|Section|Раздел)\s+(\d+|[IVXivx]+)[\.:]?\s+(.+)$',
        r'^(\d+)\.\s+([A-ZА-Я][^\n]{3,60})$',
        r'^([A-ZА-Я][^\n]{3,50})\s*$',  # all-caps or title case short lines
    ]

    toc_entries = []
    for i, line in enumerate(lines):
        line = line.strip()
        for pat in toc_patterns[:2]:
            m = re.match(pat, line, re.IGNORECASE)
            if m:
                toc_entries.append({'title': line, 'position': i})
                break

    # If no TOC found, try splitting by heading-like patterns
    if not toc_entries:
        heading_pat = re.compile(r'^(#{1,3}\s+.+|Chapter\s+\d+|Глава\s+\d+|\d+\.\s+[A-ZА-Я].{3,50})$', re.IGNORECASE)
        for i, line in enumerate(lines):
            if heading_pat.match(line.strip()) and len(line.strip()) < 80:
                toc_entries.append({'title': line.strip(), 'position': i})

    # 2. Split content by chapters
    if len(toc_entries) >= 2:
        for idx, entry in enumerate(toc_entries):
            start = entry['position']
            end = toc_entries[idx + 1]['position'] if idx + 1 < len(toc_entries) else len(lines)
            chapter_content = '\n'.join(lines[start:end]).strip()
            chapters.append({
                'title': entry['title'],
                'content': chapter_content[:10000],  # limit per chapter
                'summary': '',
                'position': start,
            })
    else:
        # Fall back to splitting into ~equal chunks of ~3000 words
        words = content.split()
        chunk_size = max(1000, len(words) // 8)
        for i in range(0, len(words), chunk_size):
            chunk_words = words[i:i + chunk_size]
            chapters.append({
                'title': f'Section {i // chunk_size + 1}',
                'content': ' '.join(chunk_words),
                'summary': '',
                'position': i,
            })

    # 3. Generate chapter summaries (limit to first 10 chapters to avoid quota)
    for chapter in chapters[:10]:
        try:
            prompt = f"""{lang_instruction}
Write a concise 2-3 sentence summary of this chapter titled "{chapter['title']}":

{chapter['content'][:5000]}

Return only the summary text, no additional formatting."""
            chapter['summary'] = _gemini(prompt)
        except Exception as e:
            chapter['summary'] = f"Summary unavailable: {str(e)[:100]}"

    # 4. Generate overall book summary
    chapter_summaries = '\n'.join(
        f"- {c['title']}: {c['summary']}" for c in chapters[:10]
    )
    try:
        prompt = f"""{lang_instruction}
Based on these chapter summaries of the book "{source.get('name', 'Book')}",
write a comprehensive 3-4 paragraph overall summary:

{chapter_summaries}

Return only the summary text."""
        overall_summary = _gemini(prompt)
    except Exception as e:
        overall_summary = f"Overall summary unavailable: {str(e)[:100]}"

    return {
        'chapters': chapters,
        'overall_summary': overall_summary,
    }


# ──────────────────────────────────────────────────────────────────────────────
# Semantic links between sources
# ──────────────────────────────────────────────────────────────────────────────

def _chunk_text(content: str, chunk_size: int = 1000) -> list:
    """Split text into chunks of approximately chunk_size words."""
    words = content.split()
    chunks = []
    for i in range(0, len(words), chunk_size):
        chunks.append(' '.join(words[i:i + chunk_size]))
    return chunks


def find_semantic_links(sources: list) -> list:
    """Find semantic connections between sources using Gemini."""
    if len(sources) < 2:
        return []

    links = []
    # Build excerpts from each source (limit to avoid huge prompts)
    source_excerpts = []
    for s in sources[:5]:  # max 5 sources
        content = s.get('content', '')[:8000]
        source_excerpts.append({
            'id': s['id'],
            'name': s.get('name', 'Unknown'),
            'excerpt': content[:3000],
        })

    # Compare pairs
    for i in range(len(source_excerpts)):
        for j in range(i + 1, len(source_excerpts)):
            sa = source_excerpts[i]
            sb = source_excerpts[j]
            try:
                prompt = f"""Analyze these two document excerpts and find semantic connections between them.

Source A: "{sa['name']}"
{sa['excerpt'][:2000]}

Source B: "{sb['name']}"
{sb['excerpt'][:2000]}

Find 2-4 specific semantic connections. For each connection return a JSON object with:
- chunk_a: relevant excerpt from Source A (max 200 chars)
- chunk_b: relevant excerpt from Source B (max 200 chars)  
- link_type: one of "confirms", "contradicts", "develops", "illustrates"
- confidence: float 0.0-1.0
- explanation: 1-2 sentence explanation in Russian

Return ONLY a valid JSON array, no markdown fences."""
                raw = _gemini(prompt)
                raw = re.sub(r'^```(?:json)?\s*', '', raw, flags=re.MULTILINE)
                raw = re.sub(r'```\s*$', '', raw, flags=re.MULTILINE)
                found_links = json.loads(raw.strip())
                if isinstance(found_links, list):
                    for link in found_links:
                        link['source_a_id'] = sa['id']
                        link['source_a_name'] = sa['name']
                        link['source_b_id'] = sb['id']
                        link['source_b_name'] = sb['name']
                        links.append(link)
            except Exception as e:
                import logging
                logging.warning(f"Semantic link analysis failed for {sa['id']}/{sb['id']}: {e}")

    return links


# ──────────────────────────────────────────────────────────────────────────────
# Brandbook processing
# ──────────────────────────────────────────────────────────────────────────────

def extract_brand_profile(content: str) -> dict:
    """Use Gemini to extract brand profile from brandbook content."""
    prompt = f"""Analyze this brandbook document and extract the brand profile.

Content:
{content[:15000]}

Return ONLY a valid JSON object (no markdown fences) with these fields:
{{
  "brand_name": "string",
  "colors": {{"primary": "hex or name", "secondary": "...", "accent": "..."}},
  "fonts": {{"primary": "font name", "secondary": "font name"}},
  "tone_of_voice": "description of brand voice and personality",
  "layout_rules": ["rule 1", "rule 2", "..."],
  "do_not_use": ["prohibited element 1", "..."],
  "tagline": "brand tagline if present",
  "target_audience": "description",
  "brand_values": ["value 1", "value 2", "..."]
}}"""
    try:
        return _gemini_json(prompt)
    except Exception as e:
        return {
            "brand_name": "Unknown",
            "colors": {},
            "fonts": {},
            "tone_of_voice": "Not extracted",
            "layout_rules": [],
            "do_not_use": [],
            "tagline": "",
            "target_audience": "",
            "brand_values": [],
            "_error": str(e),
        }


# ──────────────────────────────────────────────────────────────────────────────
# Podcast engine
# ──────────────────────────────────────────────────────────────────────────────

def gen_podcast_full(source: dict, lang: str = 'ru') -> str:
    """Generate two-host dialogue podcast (NotebookLM style, but more alive). Returns mp3 path."""
    deep     = source.get('deep_analysis') or {}
    analysis = source.get('analysis') or {}
    name     = source.get('name', 'this document')
    sid      = source.get('id', 'unknown')
    tts_lang = "ru" if lang == "ru" else "en"
    is_ru    = lang == "ru"

    host_a   = "Ведущий 1" if is_ru else "Host 1"
    host_b   = "Ведущий 2" if is_ru else "Host 2"
    lang_line = (
        "Весь диалог — ТОЛЬКО НА РУССКОМ ЯЗЫКЕ. Ни слова по-английски."
        if is_ru else
        "Write the entire transcript in ENGLISH only."
    )

    core_thesis    = deep.get('core_thesis', '')
    key_insights   = json.dumps(deep.get('key_insights',   [])[:6], ensure_ascii=False)
    logical_chains = json.dumps(deep.get('logical_chains', [])[:4], ensure_ascii=False)
    talking_points = json.dumps(deep.get('podcast_talking_points', [])[:8], ensure_ascii=False)
    practical      = json.dumps(deep.get('practical_applications', [])[:5], ensure_ascii=False)
    contradictions = json.dumps(deep.get('contradictions', [])[:3], ensure_ascii=False)
    summary        = analysis.get('summary', '')
    key_points     = ', '.join(analysis.get('key_points', [])[:6])

    prompt = f"""{lang_line}

You are writing a podcast transcript for two hosts who have ACTUALLY READ and DEEPLY UNDERSTOOD "{name}".
This is not a summary reading — it's two smart people thinking out loud together.

Hosts: {host_a} and {host_b}

Material to work from:
- Core thesis: {core_thesis}
- Key insights: {key_insights}
- Logical chains: {logical_chains}
- Talking points: {talking_points}
- Practical applications: {practical}
- Contradictions / tensions: {contradictions}
- Summary: {summary}
- Key points: {key_points}

CRITICAL STYLE RULES — READ CAREFULLY:

ENERGY & PACE:
- This is a HIGH-ENERGY conversation. Both hosts are genuinely excited, curious, or provoked.
- Short punchy exchanges. Fast back-and-forth. Don't let one person talk for more than 3 sentences.
- Frequent interruptions, incomplete sentences cut off by the other: "Wait, but—" / "No, exactly! And—"
- Enthusiasm is contagious — when one gets excited, the other catches it or pushes back harder.

VOICE & CHARACTER:
- {"Один ведущий — энтузиаст, говорит быстро, восклицает, перебивает. Второй — скептик, приземляет, говорит 'а зачем это вообще?' и 'ну это же очевидно, нет?'" if is_ru else "One host is an enthusiast — fast, excited, connecting dots everywhere. The other is a skeptic — grounding, questioning, 'but does this actually matter?'"}
- {"Разговорный живой русский: 'это же вообще...', 'подожди подожди', 'нет ну серьёзно', 'это как будто', 'слушай а ведь', 'погоди а это значит что', 'блин', 'ну вот именно!'" if is_ru else "Casual English: 'okay okay okay', 'wait that's actually—', 'no but hang on', 'right?? like—', 'I mean come on', 'that's the part that gets me'"}
- NEVER say "great point", "absolutely", "exactly right", "indeed" — these are robot words.

HUMOR:
- Jokes come from the MATERIAL ITSELF — absurd implications, unexpected connections, contradictions.
- Self-deprecating moments: "I'm definitely going to misexplain this but—"
- Exaggeration for effect: "So basically what they're saying is [hilariously oversimplified version]?"
- One host should make at least 2-3 genuinely funny observations (not forced).

SUBSTANCE:
- Use EXACT quotes, numbers, names from the material — not vague paraphrasing.
- At least one "wait so that means..." chain of logic that surprises even the host saying it.
- They should DISAGREE on at least one interpretation and NOT resolve it.
- 2000-2500 words MINIMUM. No cutting corners.

FORBIDDEN:
- No "welcome to our podcast"
- No "today we're going to discuss"  
- No "as we mentioned earlier"
- No section breaks or headers
- No summarizing what was just said

Structure (15–20 min read):

COLD OPEN — MAKE IT POP (no greeting whatsoever):
Start MID-THOUGHT. One host drops the most provocative, weird, or counterintuitive thing from the material.
It should feel like we caught them mid-conversation.
{"Example: 'Ведущий 1: Слушай, вот ты знаешь что меня реально убило в этом? [provocative fact]. Ведущий 2: Подожди, что??'" if is_ru else "Example: 'Host 1: Okay so the thing that actually broke my brain— [provocative fact]. Host 2: Wait, what?'"}
The opening must create immediate tension or curiosity. No setup, no context — throw us in.

BLOCK 1 — What IS this actually? (400–500 words)
They work out the core idea together. One might misframe it, the other corrects. Show the thinking, not the answer.

BLOCK 2 — The meat (600–700 words)
Deep dive into 2-3 specific insights. At least ONE moment where someone says something like:
"Wait, so they're saying [X]... but that means [Y] which is kind of insane because..."
Real tension, real back-and-forth. Use specific numbers, names, quotes from the material.

BLOCK 3 — Okay but who cares? (400–500 words)
Practical implications. One host pushes hard: "but what do you actually DO with this?"
The other plays devil's advocate. They should arrive at different conclusions.

BLOCK 4 — The thing that surprised them most (300–400 words)
Each picks one thing they didn't expect. These should be DIFFERENT things. Brief disagreement on which matters more.

CLOSING (150–200 words):
Each host gives their single biggest takeaway in ONE sentence. They don't agree. 
End on something memorable — a question, a provocative thought, NOT a summary.

Format STRICTLY as:
{host_a}: [text]
{host_b}: [text]
{host_a}: [text]
...

Each turn: 1–4 sentences. No long monologues — this is a DIALOGUE.
Total: 2000–2500 words.
Output ONLY the dialogue. No section headers, no stage directions, no parenthetical notes."""

    try:
        script = _gemini(prompt, quality='high')
    except Exception:
        pts = deep.get('podcast_talking_points', []) or analysis.get('key_points', [])
        script = (
            f"{host_a}: Итак, сегодня разбираем «{name}». {summary}\n"
            f"{host_b}: Интересно. Что самое важное?\n"
            + "\n".join(
                f"{host_a if i%2==0 else host_b}: {pt}"
                for i, pt in enumerate(pts[:6])
            )
        )

    # Save script as markdown
    script_path = OUTPUTS / f"{sid}_podcast_script.md"
    script_path.write_text(
        f"# Podcast Script: {name}\n\n"
        f"_Generated by Knowledge Studio_\n\n"
        f"---\n\n"
        f"{script}",
        encoding='utf-8'
    )
    sources_db[sid]["podcast_script_url"] = f"/outputs/{sid}_podcast_script.md"

    # Strip speaker labels for single-voice TTS
    tts_text = re.sub(
        rf'^(?:{re.escape(host_a)}|{re.escape(host_b)}):\s*',
        '', script, flags=re.MULTILINE
    ).strip()

    out = OUTPUTS / f"{sid}_podcast_full.mp3"
    gTTS(text=tts_text, lang=tts_lang, slow=False).save(str(out))
    return str(out)


def gen_podcast_chapter(chapter: dict, book_name: str, chapter_idx: int,
                        source_deep: dict = None, lang: str = 'ru') -> str:
    """Generate two-host dialogue podcast for a single book chapter. Returns mp3 path."""
    tts_lang = "ru" if lang == "ru" else "en"
    is_ru    = lang == "ru"

    host_a   = "Ведущий 1" if is_ru else "Host 1"
    host_b   = "Ведущий 2" if is_ru else "Host 2"
    lang_line = (
        "Весь диалог — ТОЛЬКО НА РУССКОМ ЯЗЫКЕ."
        if is_ru else
        "Write the entire transcript in ENGLISH only."
    )

    title    = chapter.get('title', f'Chapter {chapter_idx + 1}')
    content  = chapter.get('content', '')[:8000]
    summary  = chapter.get('summary', '')
    deep     = chapter.get('deep_analysis') or source_deep or {}

    talking_points = json.dumps(deep.get('podcast_talking_points', [])[:5], ensure_ascii=False)
    insights       = json.dumps(deep.get('key_insights',   [])[:4], ensure_ascii=False)
    core           = deep.get('core_thesis', '')

    prompt = f"""{lang_line}

Write a podcast transcript in the style of a NotebookLM Audio Overview, but more alive and with humor.

Two hosts ({host_a} and {host_b}) are discussing chapter "{title}" from "{book_name}".

Chapter material:
- Summary: {summary}
- Core idea: {core}
- Key insights: {insights}
- Talking points: {talking_points}
- Content: {content[:5000]}

Host style:
- Natural reactions: "Wait, so...", "Hold on...", "That's actually wild"
- Occasional disagreement on how to interpret the chapter's ideas
- Jokes when something is counterintuitive or absurd
- Real examples and analogies
- One host sometimes says "okay but who actually cares" — the other convinces them

Structure (10–15 min / 1500–2000 words):

COLD OPEN: provocative or surprising statement from the chapter (no "in this chapter...")

MAIN DISCUSSION (1000–1300 words):
- Unpack the key ideas with genuine back-and-forth
- At least one moment of genuine surprise or disagreement
- Examples, analogies, real-world connections

PRACTICAL TAKEAWAYS (250–350 words):
- What can you actually do with this? One host pushes for specifics.

CLOSING (200–250 words):
- Each host's one takeaway from the chapter — they don't fully agree
- Brief bridge: how this fits the bigger picture of the book

Format STRICTLY:
{host_a}: [text]
{host_b}: [text]
...

Each turn: 1–4 sentences. Total: 1500–2000 words.
Output ONLY the dialogue. No section headers, no stage directions."""

    try:
        script = _gemini(prompt)
    except Exception:
        script = (
            f"{host_a}: Разбираем главу «{title}» из «{book_name}». {summary}\n"
            f"{host_b}: Интересно. Что главное?\n"
            f"{host_a}: {core or content[:300]}"
        )

    # Save chapter script
    script_path = OUTPUTS / f"chapter_{chapter_idx}_podcast_script.md"
    script_path.write_text(
        f"# Chapter Podcast: {title}\n_From: {book_name}_\n\n---\n\n{script}",
        encoding='utf-8'
    )

    # Strip labels for TTS
    tts_text = re.sub(
        rf'^(?:{re.escape(host_a)}|{re.escape(host_b)}):\s*',
        '', script, flags=re.MULTILINE
    ).strip()

    out = OUTPUTS / f"chapter_{chapter_idx}_podcast.mp3"
    gTTS(text=tts_text, lang=tts_lang, slow=False).save(str(out))
    return str(out)


# ──────────────────────────────────────────────────────────────────────────────
# Merge sources helper
# ──────────────────────────────────────────────────────────────────────────────

def _merge_sources(sids: list[str]) -> dict:
    sources_list = [s for sid in sids if (s := sources_db.get(sid)) and s.get("status") == "ready"]
    if not sources_list:
        raise HTTPException(400, "No ready sources found")
    if len(sources_list) == 1:
        return sources_list[0]

    combined_name = " + ".join(s.get("name","Untitled") for s in sources_list[:3])
    if len(sources_list) > 3:
        combined_name += f" (+{len(sources_list)-3} more)"

    combined_content = "\n\n".join(
        f"=== {s.get('name','Source')} ===\n{s.get('content','')[:8000]}"
        for s in sources_list
    )
    merged_analysis: dict = {
        "summary": " | ".join(filter(None, ((s.get("analysis",{}) or {}).get("summary","") for s in sources_list)))[:2000],
        "key_points": [], "topics": [], "notable_quotes": [],
        "entities": {"people":[],"organizations":[],"places":[],"dates":[]},
        "sentiment": "mixed", "complexity": "intermediate",
        "word_count": sum((s.get("analysis",{}) or {}).get("word_count",0) for s in sources_list),
        "language": "English",
    }
    for s in sources_list:
        a = s.get("analysis",{}) or {}
        merged_analysis["key_points"].extend(a.get("key_points",[]))
        merged_analysis["topics"].extend(a.get("topics",[]))
        merged_analysis["notable_quotes"].extend(a.get("notable_quotes",[]))
        for f in ("people","organizations","places","dates"):
            merged_analysis["entities"][f].extend((a.get("entities") or {}).get(f,[]))

    for f in ("key_points","notable_quotes"):
        seen: set = set(); deduped = []
        for item in merged_analysis[f]:
            k = item[:60] if item else ""
            if k not in seen: seen.add(k); deduped.append(item)
        merged_analysis[f] = deduped[:20]
    merged_analysis["topics"] = merged_analysis["topics"][:12]
    for f in ("people","organizations","places","dates"):
        seen = set(); deduped = []
        for item in merged_analysis["entities"][f]:
            if item not in seen: seen.add(item); deduped.append(item)
        merged_analysis["entities"][f] = deduped[:15]

    return {"id":"merged_"+"_".join(sids[:3]),"name":combined_name,"type":"merged",
            "status":"ready","content":combined_content,"analysis":merged_analysis}


# ──────────────────────────────────────────────────────────────────────────────
# Background processing
# ──────────────────────────────────────────────────────────────────────────────

async def process_source(sid: str, file_path: Optional[Path], url: Optional[str], lang: str = 'ru'):
    sources_db[sid]["status"] = "processing"
    _save_sources()
    try:
        if url:
            content = await fetch_url(url)
        elif file_path:
            ext = file_path.suffix.lower()
            if ext == ".pdf":
                content = extract_pdf(file_path)
            elif ext in (".docx", ".doc"):
                content = extract_docx(file_path)
            else:
                content = extract_txt(file_path)
        else:
            content = ""

        sources_db[sid]["content_preview"] = content[:600]
        sources_db[sid]["content"] = content
        sources_db[sid]["analysis"] = analyze(content, sources_db[sid]["name"], lang)
        sources_db[sid]["status"] = "ready"
        _save_sources()

        # Deep analysis — runs automatically after regular analysis
        try:
            sources_db[sid]["deep_analysis"] = deep_analyze(sources_db[sid], lang)
        except Exception as da_err:
            import logging
            logging.warning(f"Deep analysis failed for {sid}: {da_err}")
        _save_sources()

        # Handle special source types
        source_type = sources_db[sid].get("source_type", "document")
        if source_type == "book":
            try:
                book_result = process_book(sources_db[sid], content, lang)
                sources_db[sid]["chapters"] = book_result.get("chapters", [])
                sources_db[sid]["analysis"]["summary"] = book_result.get(
                    "overall_summary", sources_db[sid]["analysis"].get("summary", "")
                )
            except Exception as book_err:
                import logging
                logging.warning(f"Book processing failed for {sid}: {book_err}")
        elif source_type == "brandbook":
            try:
                sources_db[sid]["brand_profile"] = extract_brand_profile(content)
            except Exception as brand_err:
                import logging
                logging.warning(f"Brandbook processing failed for {sid}: {brand_err}")
    except Exception as e:
        sources_db[sid]["status"] = "error"
        sources_db[sid]["error"] = str(e)
    _save_sources()


# ──────────────────────────────────────────────────────────────────────────────
# FastAPI app
# ──────────────────────────────────────────────────────────────────────────────

app = FastAPI(title="Knowledge Studio", version="4.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], allow_methods=["*"], allow_headers=["*"],
)

app.mount("/static",  StaticFiles(directory=str(STATIC)),  name="static")
app.mount("/outputs", StaticFiles(directory=str(OUTPUTS)), name="outputs")


# ──────────────────────────────────────────────────────────────────────────────
# Static routes
# ──────────────────────────────────────────────────────────────────────────────

@app.get("/")
async def index():
    return FileResponse(str(STATIC / "index.html"))

@app.get("/share/{token}")
async def share_page(token: str):
    """Serve SPA for public share view."""
    return FileResponse(str(STATIC / "index.html"))


# ──────────────────────────────────────────────────────────────────────────────
# Notebook CRUD
# ──────────────────────────────────────────────────────────────────────────────

@app.get("/api/notebooks")
async def list_notebooks():
    result = []
    for nb in notebooks_db.values():
        src_count = len(nb.get("sources", []))
        result.append({**{k: v for k, v in nb.items() if k != "sources"},
                       "source_count": src_count, "sources": nb.get("sources", [])})
    # Sort by created_at desc
    result.sort(key=lambda x: x.get("created_at",""), reverse=True)
    return result

@app.post("/api/notebooks")
async def create_notebook(payload: dict):
    nid = "nb_" + uuid.uuid4().hex[:10]
    name = payload.get("name", "").strip() or "New Notebook"
    description = payload.get("description", "").strip()
    nb = {
        "id": nid,
        "name": name,
        "description": description,
        "created_at": datetime.now().isoformat(),
        "sources": [],
        "share_token": None,
        "shared_at": None,
    }
    notebooks_db[nid] = nb
    _save_notebooks()
    return {**nb, "source_count": 0}

@app.get("/api/notebooks/{nid}")
async def get_notebook(nid: str):
    nb = notebooks_db.get(nid)
    if not nb:
        raise HTTPException(404, "Notebook not found")
    return {**nb, "source_count": len(nb.get("sources", []))}

@app.put("/api/notebooks/{nid}")
async def update_notebook(nid: str, payload: dict):
    nb = notebooks_db.get(nid)
    if not nb:
        raise HTTPException(404, "Notebook not found")
    if "name" in payload:
        nb["name"] = payload["name"].strip() or nb["name"]
    if "description" in payload:
        nb["description"] = payload["description"].strip()
    _save_notebooks()
    return {**nb, "source_count": len(nb.get("sources", []))}

@app.delete("/api/notebooks/{nid}")
async def delete_notebook(nid: str):
    nb = notebooks_db.pop(nid, None)
    if not nb:
        raise HTTPException(404, "Notebook not found")
    # Delete all sources in notebook
    for sid in nb.get("sources", []):
        s = sources_db.pop(sid, None)
        if s:
            fp = s.get("file_path")
            if fp and Path(fp).exists():
                Path(fp).unlink(missing_ok=True)
    _save_all()
    return {"deleted": nid}


# ──────────────────────────────────────────────────────────────────────────────
# Source CRUD (notebook-scoped)
# ──────────────────────────────────────────────────────────────────────────────

def _get_notebook(nid: str) -> dict:
    nb = notebooks_db.get(nid)
    if not nb:
        raise HTTPException(404, "Notebook not found")
    return nb

@app.get("/api/notebooks/{nid}/sources")
async def list_notebook_sources(nid: str):
    nb = _get_notebook(nid)
    result = []
    for sid in nb.get("sources", []):
        s = sources_db.get(sid)
        if s:
            result.append({k: v for k, v in s.items() if k != "content"})
    return result

@app.post("/api/notebooks/{nid}/sources/upload")
async def upload_to_notebook(
    nid: str,
    background_tasks: BackgroundTasks,
    file: Optional[UploadFile] = File(None),
    url: Optional[str] = Form(None),
    lang: str = Form('ru'),
    source_type: str = Form('document'),
):
    nb = _get_notebook(nid)
    sid = uuid.uuid4().hex[:10]
    try:
        if file and file.filename:
            safe = Path(file.filename).name
            if not safe or safe.startswith('.'):
                raise HTTPException(400, "Invalid filename")
            fp = UPLOADS / f"{sid}_{safe}"
            content = await file.read()
            if len(content) > 50_000_000:
                raise HTTPException(400, "File too large (max 50MB)")
            fp.write_bytes(content)
            sources_db[sid] = dict(id=sid, name=safe, type="file",
                                   source_type=source_type,
                                   file_path=str(fp), status="pending",
                                   notebook_id=nid,
                                   created_at=datetime.now().isoformat())
            nb["sources"].append(sid)
            _save_all()
            background_tasks.add_task(process_source, sid, fp, None, lang)
            return {"id": sid, "status": "pending", "name": safe}
        elif url:
            if not url.startswith(('http://','https://')):
                raise HTTPException(400, "URL must start with http:// or https://")
            sources_db[sid] = dict(id=sid, name=url[:80], type="url",
                                   source_type=source_type,
                                   url=url, status="pending",
                                   notebook_id=nid,
                                   created_at=datetime.now().isoformat())
            nb["sources"].append(sid)
            _save_all()
            background_tasks.add_task(process_source, sid, None, url, lang)
            return {"id": sid, "status": "pending", "name": url[:80]}
        else:
            raise HTTPException(400, "Provide a file or URL")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Upload failed: {str(e)}")

@app.get("/api/notebooks/{nid}/sources/{sid}")
async def get_notebook_source(nid: str, sid: str):
    nb = _get_notebook(nid)
    if sid not in nb.get("sources", []):
        raise HTTPException(404, "Source not in notebook")
    s = sources_db.get(sid)
    if not s:
        raise HTTPException(404, "Source not found")
    return {k: v for k, v in s.items() if k != "content"}

@app.delete("/api/notebooks/{nid}/sources/{sid}")
async def delete_notebook_source(nid: str, sid: str):
    nb = _get_notebook(nid)
    if sid not in nb.get("sources", []):
        raise HTTPException(404, "Source not in notebook")
    nb["sources"].remove(sid)
    s = sources_db.pop(sid, None)
    _save_all()
    if s:
        fp = s.get("file_path")
        if fp and Path(fp).exists():
            Path(fp).unlink(missing_ok=True)
    return {"deleted": sid}


# ──────────────────────────────────────────────────────────────────────────────
# Share endpoints
# ──────────────────────────────────────────────────────────────────────────────

@app.post("/api/notebooks/{nid}/share")
async def share_notebook(nid: str, request=None):
    nb = _get_notebook(nid)
    if not nb.get("share_token"):
        token = "tok_" + secrets.token_urlsafe(16)
        nb["share_token"] = token
        nb["shared_at"] = datetime.now().isoformat()
        _save_notebooks()
    token = nb["share_token"]
    return {"token": token, "url": f"/share/{token}"}

@app.delete("/api/notebooks/{nid}/share")
async def revoke_share(nid: str):
    nb = _get_notebook(nid)
    nb["share_token"] = None
    nb["shared_at"] = None
    _save_notebooks()
    return {"revoked": True}

@app.get("/api/share/{token}")
async def get_shared_notebook(token: str):
    """Public readonly view of a shared notebook."""
    nb = next((n for n in notebooks_db.values() if n.get("share_token") == token), None)
    if not nb:
        raise HTTPException(404, "Share link not found or revoked")
    sources = []
    for sid in nb.get("sources", []):
        s = sources_db.get(sid)
        if s:
            sources.append({k: v for k, v in s.items() if k not in ("content", "file_path")})
    return {
        "id": nb["id"],
        "name": nb["name"],
        "description": nb["description"],
        "created_at": nb["created_at"],
        "sources": sources,
        "token": token,
    }


# ──────────────────────────────────────────────────────────────────────────────
# Generation endpoints (notebook-scoped single source)
# ──────────────────────────────────────────────────────────────────────────────

def _require_ready_in_notebook(nid: str, sid: str) -> dict:
    nb = _get_notebook(nid)
    if sid not in nb.get("sources", []):
        raise HTTPException(404, "Source not in notebook")
    s = sources_db.get(sid)
    if not s:
        raise HTTPException(404, "Source not found")
    if s.get("status") != "ready":
        raise HTTPException(400, f"Source status is '{s.get('status')}', not ready")
    return s


@app.post("/api/notebooks/{nid}/sources/{sid}/generate/presentation")
async def nb_gen_pptx(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        out = OUTPUTS / f"{sid}_presentation.pptx"
        await asyncio.get_event_loop().run_in_executor(None, gen_presentation, s, out, lang)
        url = f"/outputs/{sid}_presentation.pptx"
        sources_db[sid]["pptx_url"] = url; _save_sources()
        return {"url": url}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/audio")
async def nb_gen_audio(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        out = OUTPUTS / f"{sid}_audio.mp3"
        await asyncio.get_event_loop().run_in_executor(None, gen_audio, s, out, lang)
        url = f"/outputs/{sid}_audio.mp3"
        sources_db[sid]["audio_url"] = url; _save_sources()
        return {"url": url}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/infographic")
async def nb_gen_infographic(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        out = OUTPUTS / f"{sid}_infographic.png"
        await asyncio.get_event_loop().run_in_executor(None, gen_infographic, s, out, lang)
        url = f"/outputs/{sid}_infographic.png"
        sources_db[sid]["infographic_url"] = url; _save_sources()
        return {"url": url}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/mindmap")
async def nb_gen_mindmap(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        out = OUTPUTS / f"{sid}_mindmap.png"
        await asyncio.get_event_loop().run_in_executor(None, gen_mindmap, s, out)
        url = f"/outputs/{sid}_mindmap.png"
        sources_db[sid]["mindmap_url"] = url; _save_sources()
        return {"url": url}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/studyguide")
async def nb_gen_study(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        text = await asyncio.get_event_loop().run_in_executor(None, gen_study_guide, s, lang)
        sources_db[sid]["study_guide"] = text; _save_sources()
        return {"content": text}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/faq")
async def nb_gen_faq(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        text = await asyncio.get_event_loop().run_in_executor(None, gen_faq, s, lang)
        sources_db[sid]["faq"] = text; _save_sources()
        return {"content": text}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/briefing")
async def nb_gen_briefing(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        text = await asyncio.get_event_loop().run_in_executor(None, gen_briefing, s, lang)
        sources_db[sid]["briefing"] = text; _save_sources()
        return {"content": text}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/timeline")
async def nb_gen_timeline(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        text = await asyncio.get_event_loop().run_in_executor(None, gen_timeline, s, lang)
        sources_db[sid]["timeline"] = text; _save_sources()
        return {"content": text}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))

@app.post("/api/notebooks/{nid}/sources/{sid}/generate/glossary")
async def nb_gen_glossary(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        text = await asyncio.get_event_loop().run_in_executor(None, gen_glossary, s, lang)
        sources_db[sid]["glossary"] = text; _save_sources()
        return {"content": text}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))


# ── New KnowledgeStudio endpoints ────────────────────────────────────────────

@app.post("/api/notebooks/{nid}/sources/{sid}/set_type")
async def set_source_type(nid: str, sid: str, payload: dict, background_tasks: BackgroundTasks):
    """Set source type and optionally trigger reprocessing."""
    nb = _get_notebook(nid)
    if sid not in nb.get("sources", []):
        raise HTTPException(404, "Source not in notebook")
    s = sources_db.get(sid)
    if not s:
        raise HTTPException(404, "Source not found")

    source_type = payload.get("type", "document")
    if source_type not in ("document", "book", "brandbook", "url", "youtube"):
        raise HTTPException(400, f"Invalid type: {source_type}")

    s["source_type"] = source_type
    lang = payload.get("lang", "ru")

    if source_type == "book" and s.get("status") == "ready":
        content = s.get("content", "")
        if content:
            background_tasks.add_task(_reprocess_book, sid, content, lang)

    elif source_type == "brandbook" and s.get("status") == "ready":
        content = s.get("content", "")
        if content:
            background_tasks.add_task(_reprocess_brandbook, sid, content)

    _save_sources()
    return {**{k: v for k, v in s.items() if k != "content"}}


async def _reprocess_book(sid: str, content: str, lang: str):
    try:
        s = sources_db.get(sid)
        if not s:
            return
        result = await asyncio.get_event_loop().run_in_executor(
            None, process_book, s, content, lang
        )
        s["chapters"] = result.get("chapters", [])
        s["analysis"] = s.get("analysis") or {}
        s["analysis"]["summary"] = result.get("overall_summary", s["analysis"].get("summary", ""))
        _save_sources()
    except Exception as e:
        import logging
        logging.error(f"Book reprocess failed for {sid}: {e}")


async def _reprocess_brandbook(sid: str, content: str):
    try:
        s = sources_db.get(sid)
        if not s:
            return
        profile = await asyncio.get_event_loop().run_in_executor(
            None, extract_brand_profile, content
        )
        s["brand_profile"] = profile
        _save_sources()
    except Exception as e:
        import logging
        logging.error(f"Brandbook reprocess failed for {sid}: {e}")


@app.get("/api/notebooks/{nid}/links")
async def get_notebook_links(nid: str):
    """Return semantic links between sources in this notebook."""
    nb = _get_notebook(nid)
    links = nb.get("semantic_links", [])
    return {"links": links}


@app.post("/api/notebooks/{nid}/analyze_links")
async def analyze_notebook_links(nid: str, background_tasks: BackgroundTasks):
    """Trigger semantic link analysis between sources (async)."""
    nb = _get_notebook(nid)
    ready_sources = [
        sources_db[sid] for sid in nb.get("sources", [])
        if sources_db.get(sid) and sources_db[sid].get("status") == "ready"
    ]
    if len(ready_sources) < 2:
        raise HTTPException(400, "Need at least 2 ready sources to analyze links")
    nb["links_status"] = "analyzing"
    _save_notebooks()
    background_tasks.add_task(_run_link_analysis, nid, ready_sources)
    return {"status": "analyzing", "source_count": len(ready_sources)}


async def _run_link_analysis(nid: str, ready_sources: list):
    try:
        links = await asyncio.get_event_loop().run_in_executor(
            None, find_semantic_links, ready_sources
        )
        nb = notebooks_db.get(nid)
        if nb:
            nb["semantic_links"] = links
            nb["links_status"] = "ready"
            nb["links_analyzed_at"] = datetime.now().isoformat()
            _save_notebooks()
    except Exception as e:
        nb = notebooks_db.get(nid)
        if nb:
            nb["links_status"] = "error"
            nb["links_error"] = str(e)
            _save_notebooks()


@app.post("/api/notebooks/{nid}/sources/{sid}/generate/podcast_full")
async def nb_gen_podcast_full(nid: str, sid: str, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        out_path = await asyncio.get_event_loop().run_in_executor(
            None, gen_podcast_full, s, lang
        )
        url = f"/outputs/{Path(out_path).name}"
        sources_db[sid]["podcast_full_url"] = url
        _save_sources()
        return {"url": url}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/api/notebooks/{nid}/sources/{sid}/generate/podcast_chapter/{chapter_idx}")
async def nb_gen_podcast_chapter(nid: str, sid: str, chapter_idx: int, lang: str = Query('ru')):
    try:
        s = _require_ready_in_notebook(nid, sid)
        chapters = s.get("chapters", [])
        if not chapters:
            raise HTTPException(400, "No chapters found. Make sure source is processed as a book.")
        if chapter_idx < 0 or chapter_idx >= len(chapters):
            raise HTTPException(400, f"Chapter index {chapter_idx} out of range (0-{len(chapters)-1})")
        chapter = chapters[chapter_idx]
        book_name = s.get("name", "Book")
        source_deep = s.get("deep_analysis") or {}
        out_path = await asyncio.get_event_loop().run_in_executor(
            None, gen_podcast_chapter, chapter, book_name, chapter_idx, source_deep, lang
        )
        url = f"/outputs/{Path(out_path).name}"
        return {"url": url, "chapter_title": chapter.get("title", "")}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, str(e))


@app.get("/api/notebooks/{nid}/sources/{sid}/podcast_script")
async def nb_get_podcast_script(nid: str, sid: str):
    s = _require_ready_in_notebook(nid, sid)
    script_url = s.get("podcast_script_url")
    if not script_url:
        raise HTTPException(404, "Podcast script not yet generated")
    script_path = OUTPUTS / f"{sid}_podcast_script.md"
    if not script_path.exists():
        raise HTTPException(404, "Script file not found")
    return {"content": script_path.read_text(encoding='utf-8'), "url": script_url}


@app.get("/api/notebooks/{nid}/sources/{sid}/deep_analysis")
async def nb_get_deep_analysis(nid: str, sid: str):
    s = _require_ready_in_notebook(nid, sid)
    da = s.get("deep_analysis")
    if not da:
        raise HTTPException(404, "Deep analysis not yet available for this source")
    return da


# ── Multi-source generation (notebook-scoped) ─────────────────────────────────

@app.post("/api/notebooks/{nid}/generate/{gen_type}")
async def nb_multi_generate(nid: str, gen_type: str, payload: dict):
    try:
        nb = _get_notebook(nid)
        source_ids = payload.get("source_ids", [])
        lang = payload.get("lang", "ru")
        # Validate all sids belong to notebook
        valid_sids = [sid for sid in source_ids if sid in nb.get("sources", [])]
        if not valid_sids:
            valid_sids = nb.get("sources", [])
        if not valid_sids:
            raise HTTPException(400, "No sources in notebook")

        merged = _merge_sources(valid_sids)
        key = nid + "_" + "_".join(sorted(valid_sids))[:30]

        if gen_type == "presentation":
            out = OUTPUTS / f"{key}_presentation.pptx"
            await asyncio.get_event_loop().run_in_executor(None, gen_presentation, merged, out, lang)
            return {"url": f"/outputs/{key}_presentation.pptx"}
        elif gen_type == "audio":
            out = OUTPUTS / f"{key}_audio.mp3"
            await asyncio.get_event_loop().run_in_executor(None, gen_audio, merged, out, lang)
            return {"url": f"/outputs/{key}_audio.mp3"}
        elif gen_type == "infographic":
            out = OUTPUTS / f"{key}_infographic.png"
            await asyncio.get_event_loop().run_in_executor(None, gen_infographic, merged, out, lang)
            return {"url": f"/outputs/{key}_infographic.png"}
        elif gen_type == "mindmap":
            out = OUTPUTS / f"{key}_mindmap.png"
            await asyncio.get_event_loop().run_in_executor(None, gen_mindmap, merged, out)
            return {"url": f"/outputs/{key}_mindmap.png"}
        elif gen_type == "studyguide":
            text = await asyncio.get_event_loop().run_in_executor(None, gen_study_guide, merged, lang)
            return {"content": text}
        elif gen_type == "faq":
            text = await asyncio.get_event_loop().run_in_executor(None, gen_faq, merged, lang)
            return {"content": text}
        elif gen_type == "briefing":
            text = await asyncio.get_event_loop().run_in_executor(None, gen_briefing, merged, lang)
            return {"content": text}
        elif gen_type == "timeline":
            text = await asyncio.get_event_loop().run_in_executor(None, gen_timeline, merged, lang)
            return {"content": text}
        elif gen_type == "glossary":
            text = await asyncio.get_event_loop().run_in_executor(None, gen_glossary, merged, lang)
            return {"content": text}
        else:
            raise HTTPException(400, f"Unknown generation type: {gen_type}")
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, str(e))


# ── Legacy endpoints (kept for compatibility) ─────────────────────────────────

@app.get("/api/sources")
async def list_sources_legacy():
    return [{k: v for k, v in s.items() if k != "content"} for s in sources_db.values()]

@app.post("/api/sources/upload")
async def upload_legacy(background_tasks: BackgroundTasks,
                        file: Optional[UploadFile] = File(None),
                        url: Optional[str] = Form(None)):
    # Find or create default notebook
    if not notebooks_db:
        nid = "nb_" + uuid.uuid4().hex[:10]
        notebooks_db[nid] = {"id": nid, "name": "My Research", "description": "",
                              "created_at": datetime.now().isoformat(),
                              "sources": [], "share_token": None, "shared_at": None}
        _save_notebooks()
    nid = next(iter(notebooks_db))
    return await upload_to_notebook(nid, background_tasks, file, url)

@app.get("/api/sources/{sid}")
async def get_source_legacy(sid: str):
    s = sources_db.get(sid)
    if not s: raise HTTPException(404, "Not found")
    return {k: v for k, v in s.items() if k != "content"}

@app.delete("/api/sources/{sid}")
async def delete_source_legacy(sid: str):
    s = sources_db.pop(sid, None)
    if not s: raise HTTPException(404, "Not found")
    # Remove from notebook
    for nb in notebooks_db.values():
        if sid in nb.get("sources", []):
            nb["sources"].remove(sid)
    _save_all()
    fp = s.get("file_path")
    if fp and Path(fp).exists():
        Path(fp).unlink(missing_ok=True)
    return {"deleted": sid}

@app.post("/api/generate/{gen_type}")
async def multi_generate_legacy(gen_type: str, payload: dict):
    source_ids = payload.get("source_ids", [])
    lang = payload.get("lang", "ru")
    if not source_ids:
        raise HTTPException(400, "No source_ids provided")
    merged = _merge_sources(source_ids)
    key = "_".join(sorted(source_ids))[:40]

    if gen_type == "presentation":
        out = OUTPUTS / f"{key}_presentation.pptx"
        await asyncio.get_event_loop().run_in_executor(None, gen_presentation, merged, out, lang)
        return {"url": f"/outputs/{key}_presentation.pptx"}
    elif gen_type == "audio":
        out = OUTPUTS / f"{key}_audio.mp3"
        await asyncio.get_event_loop().run_in_executor(None, gen_audio, merged, out, lang)
        return {"url": f"/outputs/{key}_audio.mp3"}
    elif gen_type == "infographic":
        out = OUTPUTS / f"{key}_infographic.png"
        await asyncio.get_event_loop().run_in_executor(None, gen_infographic, merged, out, lang)
        return {"url": f"/outputs/{key}_infographic.png"}
    elif gen_type == "mindmap":
        out = OUTPUTS / f"{key}_mindmap.png"
        await asyncio.get_event_loop().run_in_executor(None, gen_mindmap, merged, out)
        return {"url": f"/outputs/{key}_mindmap.png"}
    elif gen_type == "studyguide":
        return {"content": await asyncio.get_event_loop().run_in_executor(None, gen_study_guide, merged, lang)}
    elif gen_type == "faq":
        return {"content": await asyncio.get_event_loop().run_in_executor(None, gen_faq, merged, lang)}
    elif gen_type == "briefing":
        return {"content": await asyncio.get_event_loop().run_in_executor(None, gen_briefing, merged, lang)}
    elif gen_type == "timeline":
        return {"content": await asyncio.get_event_loop().run_in_executor(None, gen_timeline, merged, lang)}
    elif gen_type == "glossary":
        return {"content": await asyncio.get_event_loop().run_in_executor(None, gen_glossary, merged, lang)}
    else:
        raise HTTPException(400, f"Unknown generation type: {gen_type}")


# ── Chat ──────────────────────────────────────────────────────────────────────

@app.post("/api/chat")
async def chat(payload: dict):
    try:
        question     = payload.get("question", "").strip()
        source_ids   = payload.get("source_ids", [])
        chat_history = payload.get("history", [])
        if not question:
            raise HTTPException(400, "No question provided")

        context_parts = []
        for sid in source_ids:
            s = sources_db.get(sid)
            if not s or s.get("status") != "ready": continue
            content = s.get("content", ""); a = s.get("analysis", {}) or {}
            part = f"=== SOURCE: {s['name']} ===\n"
            part += content[:10000] if content else f"Summary: {a.get('summary','')}\nKey points: {'; '.join(a.get('key_points',[]))}\n"
            context_parts.append(part)

        history_text = "".join(f"\n{m.get('role','user').upper()}: {m.get('content','')}"
                               for m in chat_history[-6:])

        if not context_parts:
            prompt = f"You are NoteLM, a helpful AI research assistant.\n{history_text}\nUSER: {question}\nASSISTANT:"
        else:
            context = "\n\n".join(context_parts)
            prompt = f"""You are NoteLM, an expert research assistant. Answer questions based on the provided sources.
Be precise, cite sources by name, and be thorough. If the answer is not in the sources, say so.

SOURCES:
{context[:28000]}

CONVERSATION HISTORY:{history_text}

USER: {question}
ASSISTANT:"""

        resp = await asyncio.get_event_loop().run_in_executor(None, _gemini, prompt)
        return {"answer": resp}
    except HTTPException: raise
    except Exception as e: raise HTTPException(500, f"Chat error: {str(e)}")


# ──────────────────────────────────────────────────────────────────────────────
# Health
# ──────────────────────────────────────────────────────────────────────────────

@app.get("/api/health")
async def health():
    return {"status": "ok", "version": "4.0.0",
            "notebooks": len(notebooks_db), "sources": len(sources_db)}

@app.on_event("startup")
async def startup_event():
    _load_dbs()
    _migrate()
    for sid, s in sources_db.items():
        if s.get("status") == "pending":
            s["status"] = "error"
            s["error"] = "Interrupted by restart"
    _save_sources()

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=80, reload=False)
